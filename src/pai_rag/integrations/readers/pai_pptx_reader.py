"""Pptx parser.

"""
from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, List, Optional, Union, Any
import os
import aspose.slides as slides
import aspose.pydrawing as drawing
from PIL import Image
import time
from pai_rag.utils.markdown_utils import (
    transform_local_to_oss,
    convert_table_to_markdown,
    PaiTable,
)
from pathlib import Path
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from loguru import logger


class PaiPptxReader(BaseReader):
    def __init__(
        self,
        enable_table_summary: bool = False,
        oss_cache: Any = None,
    ) -> None:
        self.enable_table_summary = enable_table_summary
        self._oss_cache = oss_cache
        logger.info(
            f"PaiPptxReader created with enable_table_summary : {self.enable_table_summary}"
        )

    def _extract_shape(self, slide_number, shape):
        image_flag = False
        markdown = []
        if shape.name.startswith("Title"):
            # 标题
            markdown.append(f"# {shape.text}\n\n")
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # 图片
            image_flag = True
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # 文本框
            markdown.append(f"{shape.text}\n\n")
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # 表格
            table = shape.table
            markdown.append(self._convert_table_to_pai_table(table))
            markdown.append("\n\n")
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            texts = []
            for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                md, shape_image_flag = self._extract_shape(slide_number, p)
                if md:
                    texts.append(md)
                if shape_image_flag:
                    image_flag = shape_image_flag

            markdown.append("\n".join(texts))

        return "".join(markdown), image_flag

    def _convert_table_to_pai_table(self, table):
        table_matrix = [
            ["" for _ in range(len(table.columns))] for _ in range(len(table.rows))
        ]
        for i in range(len(table.rows)):
            for j in range(len(table.columns)):
                cell_content = table.cell(i, j).text.replace("\n", "").replace("\r", "")
                if table.cell(i, j).is_merge_origin:
                    col_span = table.cell(i, j).span_width
                    row_span = table.cell(i, j).span_height
                    while (
                        col_span > 1
                        and j + col_span <= len(table.columns)
                        and table_matrix[i][j + col_span - 1] == ""
                    ):
                        col_span -= 1
                        table_matrix[i][j + col_span] = cell_content
                    while (
                        row_span > 1
                        and i + row_span <= len(table.rows)
                        and table_matrix[i + row_span - 1][j] == ""
                    ):
                        row_span -= 1
                        table_matrix[i + row_span][j] = cell_content
                if table_matrix[i][j] == "":
                    table_matrix[i][j] = cell_content

        row_headers_index = []
        col_headers_index = []
        if table.first_row:
            row_headers_index.append(0)
        if table.first_col:
            col_headers_index.append(0)
        pai_table = PaiTable(
            data=table_matrix,
            row_headers_index=row_headers_index,
            column_headers_index=col_headers_index,
        )
        return convert_table_to_markdown(pai_table, len(table.columns))

    def convert_pptx_to_markdown(self, fnm):
        ppt_name = os.path.basename(fnm).split(".")[0]
        ppt_name = ppt_name.replace(" ", "_")
        fnm = fnm.as_posix()
        prs = Presentation(fnm)

        markdown = []
        slide_image_flag = []
        for slide_number, slide in enumerate(prs.slides, start=1):
            image_flag = False
            for shape in slide.shapes:
                shape_markdown, shape_image_flag = self._extract_shape(
                    slide_number, shape
                )
                if shape_image_flag:
                    image_flag = shape_image_flag
                markdown.append(shape_markdown)
            markdown.append(f"# slide_number_{slide_number}\n\n")
            slide_image_flag.append(image_flag)

        if self._oss_cache:
            with slides.Presentation(fnm) as presentation:
                for i, slide in enumerate(presentation.slides):
                    if slide_image_flag[i]:
                        buffered = BytesIO()
                        slide.get_thumbnail(0.5, 0.5).save(
                            buffered, drawing.imaging.ImageFormat.jpeg
                        )
                        buffered.seek(0)
                        image = Image.open(buffered)
                        image_url = transform_local_to_oss(
                            self._oss_cache, image, ppt_name
                        )
                        time_tag = int(time.time())
                        alt_text = f"pai_rag_image_{time_tag}_"
                        image_content = f"![{alt_text}]({image_url})"
                        markdown.append(f"{image_content}\n\n")

        return "".join(markdown)

    def load_data(
        self,
        file_path: Union[Path, str],
        metadata: bool = True,
        extra_info: Optional[Dict] = None,
    ) -> List[Document]:
        """Loads list of documents from Pptx file and also accepts extra information in dict format."""
        return self.load(file_path, metadata=metadata, extra_info=extra_info)

    def load(
        self,
        file_path: Union[Path, str],
        metadata: bool = True,
        extra_info: Optional[Dict] = None,
    ) -> List[Document]:
        """Loads list of documents from Pptx file and also accepts extra information in dict format.

        Args:
            file_path (Union[Path, str]): file path of Pptx file (accepts string or Path).
            metadata (bool, optional): if metadata to be included or not. Defaults to True.
            extra_info (Optional[Dict], optional): extra information related to each document in dict format. Defaults to None.

        Raises:
            TypeError: if extra_info is not a dictionary.
            TypeError: if file_path is not a string or Path.

        Returns:
            List[Document]: list of documents.
        """
        md_content = self.convert_pptx_to_markdown(file_path)
        logger.info(f"[PaiPptxReader] successfully processed pptx file {file_path}.")
        docs = []
        if metadata and extra_info:
            extra_info = extra_info
        else:
            extra_info = dict()
            logger.info(f"processed pptx file {file_path} without metadata")
        doc = Document(text=md_content, extra_info=extra_info)
        docs.append(doc)
        logger.info(f"[PaiPptxReader] successfully loaded {len(docs)} nodes.")
        return docs
